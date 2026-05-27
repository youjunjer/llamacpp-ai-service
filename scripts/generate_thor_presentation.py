from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = Path(r"F:\GoogleDrv\P_程式開發\202604_Ollama\docs\Jetson_Thor_邊緣AI展示簡報_v1.pptx")

TITLE_COLOR = RGBColor(21, 30, 45)
ACCENT = RGBColor(0, 176, 155)
TEXT = RGBColor(55, 65, 81)
MUTED = RGBColor(99, 115, 129)
BG = RGBColor(245, 248, 252)
WHITE = RGBColor(255, 255, 255)


SLIDES = [
    {
        "title": "邊緣 AI 與多模態展示",
        "subtitle": "Jetson Thor 平台實作\n從語音、視覺理解到即時辨識",
        "type": "cover",
    },
    {
        "title": "AI 正在從雲端走向終端設備",
        "bullets": [
            "過去多數 AI 應用依賴雲端模型與資料中心運算。",
            "現在更多情境需要即時反應、資料不出場域、甚至離線運作。",
            "因此 AI 正從『雲端推論』逐步延伸到『邊緣端推論』。",
        ],
        "footer": "關鍵字：低延遲、隱私、離線、多模態",
    },
    {
        "title": "什麼是邊緣 AI",
        "bullets": [
            "在本地設備直接執行模型推論，不必每次把資料送到雲端。",
            "可直接結合相機、麥克風、感測器與現場裝置。",
            "更適合做即時互動、設備控制與現場決策。",
        ],
        "footer": "邊緣 AI 不只是部署位置不同，而是整體互動模式不同。",
    },
    {
        "title": "為什麼邊緣 AI 很重要",
        "bullets": [
            "低延遲：不用雲端 round-trip，現場回應更快。",
            "高隱私：影像、語音、文件可留在本地處理。",
            "離線可用：弱網路或無網路環境仍可運作。",
            "成本可控：減少長期傳輸與雲端推論成本。",
        ],
        "footer": "對工廠、交通、校園、醫療與零售都特別重要。",
    },
    {
        "title": "邊緣 AI 的典型應用場景",
        "bullets": [
            "智慧巡檢：設備面板、儀表、警示燈與現場異常判讀。",
            "文件理解：發票、藥袋、表單、名片與海報摘要。",
            "智慧交通：車流、人流、路況擁擠程度分析。",
            "教育與零售：互動助理、展示裝置、智慧看板。",
            "機器人與實體 AI：直接理解感測輸入並做本地決策。",
        ],
    },
    {
        "title": "為什麼邊緣 AI 需要高效能平台",
        "bullets": [
            "多模態模型不只吃文字，還包含影像、語音與串流資料。",
            "要同時兼顧模型大小、記憶體、頻寬與 I/O 整合能力。",
            "真正可用的邊緣 AI 平台，必須從 demo 走向持續運作。",
        ],
        "footer": "因此平台能力，決定了 AI 能否真的進入產品與場域。",
    },
    {
        "title": "Jetson Thor 的定位",
        "bullets": [
            "NVIDIA 新一代高效能邊緣 AI 平台。",
            "面向機器人、Physical AI、多模態推論與高互動設備。",
            "不只是開發板，而是可承載大型模型的邊緣 AI 主機。",
        ],
        "footer": "Thor 的價值在於：把大型多模態 AI 帶到終端設備。",
    },
    {
        "title": "Jetson Thor 主要規格",
        "bullets": [
            "GPU：NVIDIA Blackwell",
            "AI 算力：最高約 2070 FP4 TFLOPS",
            "CPU：14-core Arm Neoverse-V3AE",
            "記憶體：128GB LPDDR5X",
            "記憶體頻寬：273 GB/s",
            "儲存：1TB NVMe",
            "功耗範圍：40W–130W",
        ],
        "footer": "規格重點：可承載大型多模態模型與即時任務。",
    },
    {
        "title": "把規格翻譯成真正的應用能力",
        "bullets": [
            "可在本地端執行大型多模態模型。",
            "可同時處理語音、影像、偵測等複合任務。",
            "可接相機、麥克風、影片來源與自製 WebUI。",
            "適合從展示原型延伸到設備級產品。",
        ],
    },
    {
        "title": "我們這次的展示系統",
        "bullets": [
            "WebUI 展示介面：由瀏覽器直接操作。",
            "Gemma：負責語言與視覺理解。",
            "YOLO：負責即時物件辨識與 FPS 展示。",
            "輸入來源：相機、圖片網址、語音。",
            "輸出形式：文字、語音、JSON、即時標註畫面。",
        ],
    },
    {
        "title": "視覺理解展示內容",
        "bullets": [
            "發票辨識：讀取店名、日期、金額與主要品項。",
            "藥袋辨識：整理藥名、服用方式與說明。",
            "餐點熱量分析：估算組成、熱量與營養成分。",
            "路況壅擠程度：輸出 JSON 與文字描述。",
            "人流計算：估計人數並描述場景。",
        ],
        "footer": "重點不是看圖，而是將圖像整理成可用資訊。",
    },
    {
        "title": "語音互動與即時辨識展示",
        "bullets": [
            "語音輸入：按鍵錄音、送出後由模型理解內容。",
            "串流回覆：文字逐步顯示，並可語音播報。",
            "YOLO 即時辨識：以 webcam 或影片來源展示速度與辨識能力。",
            "現場展示重點：互動性、即時性、設備端執行。",
        ],
    },
    {
        "title": "現場 Demo 流程",
        "bullets": [
            "Demo 1：語音對話開場，快速建立互動感。",
            "Demo 2：視覺理解主秀，展示文件與場景分析。",
            "Demo 3：YOLO 即時辨識，展示視覺推論速度與穩定度。",
            "建議控制：每段約 3–5 分鐘，避免卡在單一操作。",
        ],
        "footer": "整體節奏：先互動、再理解、最後用即時辨識收尾。",
    },
    {
        "title": "結語",
        "bullets": [
            "AI 正在從雲端走向終端設備。",
            "邊緣 AI 的核心價值是低延遲、隱私、離線與整合能力。",
            "Jetson Thor 讓大型多模態 AI 真正有機會在本地端落地。",
            "我們這次展示的是一套可互動、可延伸、可產品化的邊緣 AI 系統。",
        ],
    },
]


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def style_text_frame(tf, font_size=24, color=TEXT):
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = "Microsoft JhengHei"
            r.font.size = Pt(font_size)
            r.font.color.rgb = color


def add_top_bar(slide):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.18)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_title(slide, title):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.55), Inches(11.5), Inches(0.8))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Microsoft JhengHei"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.75), Inches(6.7), Inches(11.6), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft JhengHei"
    run.font.size = Pt(13)
    run.font.color.rgb = MUTED


def add_bullets(slide, bullets):
    box = slide.shapes.add_textbox(Inches(0.95), Inches(1.45), Inches(11.2), Inches(4.9))
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(14)
        p.line_spacing = 1.25
    style_text_frame(tf, font_size=23)


def add_cover(prs, slide, title, subtitle):
    set_bg(slide)
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.7), Inches(5.8)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = WHITE
    accent.line.color.rgb = ACCENT
    accent.line.width = Pt(2)

    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.8), Inches(0.8), Inches(0.24), Inches(5.8)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.35), Inches(1.5), Inches(10.2), Inches(1.4))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Microsoft JhengHei"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = TITLE_COLOR

    sub_box = slide.shapes.add_textbox(Inches(1.4), Inches(2.85), Inches(9.5), Inches(1.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    for i, line in enumerate(subtitle.splitlines()):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = line
        run.font.name = "Microsoft JhengHei"
        run.font.size = Pt(20)
        run.font.color.rgb = TEXT
        para.space_after = Pt(10)

    note_box = slide.shapes.add_textbox(Inches(1.4), Inches(5.45), Inches(8.5), Inches(0.5))
    tf = note_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "1 小時演講與展示建議版"
    r.font.name = "Microsoft JhengHei"
    r.font.size = Pt(14)
    r.font.color.rgb = MUTED


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for spec in SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        if spec.get("type") == "cover":
            add_cover(prs, slide, spec["title"], spec["subtitle"])
            continue

        set_bg(slide)
        add_top_bar(slide)
        add_title(slide, spec["title"])
        add_bullets(slide, spec["bullets"])
        if spec.get("footer"):
            add_footer(slide, spec["footer"])

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
